from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Namma Bengaluru Clean"
    subtitle.text = "AI-Powered Civic Engagement & Infrastructure Monitoring\nPowered by CivicLens AI"

    # --- Slide 2: The Problem ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem Statement"
    content = slide.placeholders[1].text_frame
    content.text = "Manual reporting is slow and unverified"
    content.add_paragraph().text = "High volume of spam/irrelevant data"
    content.add_paragraph().text = "Redundant reports for the same issue (noise)"
    content.add_paragraph().text = "Lack of real-time urgency prioritization"

    # --- Slide 3: The Solution ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The CivicLens Solution"
    content = slide.placeholders[1].text_frame
    content.text = "Zero-Shot AI Verification (NVIDIA Vision AI)"
    content.add_paragraph().text = "Spatial Duplicate Detection (PostGIS 10m Radius)"
    content.add_paragraph().text = "Gamified Rewards (Impact Points)"
    content.add_paragraph().text = "Administrative Urgency Heatmaps"

    # --- Slide 4: Tech Stack ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technology Stack"
    content = slide.placeholders[1].text_frame
    content.text = "Frontend: React.js + Tailwind CSS"
    content.add_paragraph().text = "Backend: Node.js + Express + PostGIS"
    content.add_paragraph().text = "AI Service: Python FastAPI + NVIDIA NIM"
    content.add_paragraph().text = "Cloud AI: Meta LLaMA 3.2 Vision"

    # --- Slide 5: AI Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI-Powered Verification"
    content = slide.placeholders[1].text_frame
    content.text = "Zero-Shot Prompting: No custom training needed"
    content.add_paragraph().text = "Strict Anti-Spam: Blocks dogs, chairs, and indoor photos"
    content.add_paragraph().text = "Heuristic Logic: High-confidence potholes verified instantly"
    content.add_paragraph().text = "Fallback: Low-confidence reports sent for manual review"

    # --- Slide 6: Gamification ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Community Impact Points"
    content = slide.placeholders[1].text_frame
    content.text = "New Valid Report: +10 Points"
    content.add_paragraph().text = "Upvoting Duplicate: +5 Points"
    content.add_paragraph().text = "Impact: Crowdsourced confirmation increases priority score"

    # --- Slide 7: Admin Dashboard ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Administrative Spatial Analytics"
    content = slide.placeholders[1].text_frame
    content.text = "Real-time Urgency Heatmap (Heat-layers)"
    content.add_paragraph().text = "Ward-level resolution tracking"
    content.add_paragraph().text = "Automated Clean-up Drive scheduling"

    # --- Slide 8: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Impact & Future Roadmap"
    content = slide.placeholders[1].text_frame
    content.text = "90% reduction in manual verification time"
    content.add_paragraph().text = "Cleaner urban data for better planning"
    content.add_paragraph().text = "Scaleable to streetlights, encroachments, etc."

    # Save
    prs.save("CivicLens_AI_Presentation.pptx")
    print("Presentation saved successfully!")

if __name__ == "__main__":
    create_presentation()
